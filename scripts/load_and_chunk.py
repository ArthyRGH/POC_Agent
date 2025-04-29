import os
import sys
from pathlib import Path
import re
import json
from typing import List, Dict, Any
import pandas as pd
from ebooklib import epub
from lxml import etree
from unstructured.partition.auto import partition
from unstructured.chunking.title import chunk_by_title
import time
import subprocess
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

# Import document processing libraries
try:
    import pypdf
    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False
    print("Warning: pypdf package not available. PDF processing will be limited.")
    print("To install: pip install pypdf")

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Warning: python-docx package not available. DOCX processing will be limited.")
    print("To install: pip install python-docx")

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    print("Warning: openpyxl package not available. Excel processing will be limited.")
    print("To install: pip install openpyxl")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx package not available. PPTX processing will be limited.")
    print("To install: pip install python-pptx")

try:
    import rtf
    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False
    print("Warning: python-rtf package not available. RTF processing will be limited.")
    print("To install: pip install python-rtf")

# List of supported image extensions (for filtering)
IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'}

def clean_text(text: str) -> str:
    """Clean text by removing extra whitespace and non-printable characters"""
    if not text:
        return ""

    # Replace multiple whitespace with a single space
    text = re.sub(r'\s+', ' ', text)

    # Remove non-printable characters
    text = ''.join(c if ord(c) >= 32 or c in '\n\r\t' else ' ' for c in text)

    return text.strip()

def process_text_file(filepath: str) -> List[str]:
    """Process a plain text file"""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
            content = clean_text(content)
            return chunk_text(content)
    except Exception as e:
        print(f"Error processing text file {filepath}: {e}")
        return []

def process_pdf_file(filepath: str) -> List[str]:
    """Process a PDF file"""
    if PYPDF_AVAILABLE:
        chunks = []
        try:
            with open(filepath, 'rb') as f:
                pdf = pypdf.PdfReader(f)
                full_text = ""

                for page_num in range(len(pdf.pages)):
                    page = pdf.pages[page_num]
                    text = page.extract_text()
                    if text:
                        full_text += text + "\n\n"

                full_text = clean_text(full_text)
                return chunk_text(full_text)
        except Exception as e:
            print(f"Error processing PDF {filepath}: {e}")
            return []
    return ["[PDF content not extracted - pypdf not available]"]

def process_docx_file(filepath: str) -> List[str]:
    """Process a Word document"""
    if DOCX_AVAILABLE:
        try:
            doc = docx.Document(filepath)
            full_text = ""

            for para in doc.paragraphs:
                if para.text:
                    full_text += para.text + "\n\n"

            full_text = clean_text(full_text)
            return chunk_text(full_text)
        except Exception as e:
            print(f"Error processing DOCX {filepath}: {e}")
            return []
    return ["[DOCX content not extracted - python-docx not available]"]

def process_excel_file(filepath: str) -> List[str]:
    """Process an Excel file"""
    if EXCEL_AVAILABLE:
        try:
            df = pd.read_excel(filepath)
            full_text = ""
            
            # Convert each sheet to text
            for sheet_name in df.sheet_names:
                sheet_df = pd.read_excel(filepath, sheet_name=sheet_name)
                full_text += f"Sheet: {sheet_name}\n\n"
                full_text += sheet_df.to_string() + "\n\n"
            
            full_text = clean_text(full_text)
            return chunk_text(full_text)
        except Exception as e:
            print(f"Error processing Excel file {filepath}: {e}")
            return []
    return ["[Excel content not extracted - openpyxl not available]"]

def process_epub_file(filepath: str) -> List[str]:
    """Process an EPUB file"""
    try:
        book = epub.read_epub(filepath)
        full_text = ""
        
        for item in book.get_items():
            if item.get_type() == epub.ITEM_DOCUMENT:
                full_text += item.get_content().decode('utf-8') + "\n\n"
        
        full_text = clean_text(full_text)
        return chunk_text(full_text)
    except Exception as e:
        print(f"Error processing EPUB {filepath}: {e}")
        return []

def process_pptx_file(filepath: str) -> List[str]:
    """Process a PowerPoint file"""
    if PPTX_AVAILABLE:
        try:
            prs = Presentation(filepath)
            full_text = ""
            
            for slide in prs.slides:
                # Add slide title if available
                if slide.shapes.title:
                    full_text += slide.shapes.title.text + "\n\n"
                
                # Process all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        full_text += shape.text.strip() + "\n\n"
            
            full_text = clean_text(full_text)
            return chunk_text(full_text)
        except Exception as e:
            print(f"Error processing PPTX {filepath}: {e}")
            return []
    return ["[PPTX content not extracted - python-pptx not available]"]

def process_rtf_file(filepath: str) -> List[str]:
    """Process an RTF file"""
    if RTF_AVAILABLE:
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
                text = rtf.Rtf15Reader.read(rtf_content)
                text = clean_text(text)
                return chunk_text(text)
        except Exception as e:
            print(f"Error processing RTF {filepath}: {e}")
            return []
    return ["[RTF content not extracted - python-rtf not available]"]

def process_json_file(filepath: str) -> List[str]:
    """Process a JSON file"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            data = json.load(f)
            text = json.dumps(data, indent=2)
            text = clean_text(text)
            return chunk_text(text)
    except Exception as e:
        print(f"Error processing JSON {filepath}: {e}")
        return []

def process_xml_file(filepath: str) -> List[str]:
    """Process an XML file"""
    try:
        tree = etree.parse(filepath)
        text = etree.tostring(tree, pretty_print=True, encoding='unicode')
        text = clean_text(text)
        return chunk_text(text)
    except Exception as e:
        print(f"Error processing XML {filepath}: {e}")
        return []

def chunk_text(text: str, chunk_size: int = 500, chunk_overlap: int = 100) -> List[str]:
    """Split text into chunks with overlap"""
    if not text:
        return []
    
    chunks = []
    current_position = 0
    
    while current_position < len(text):
        chunk_end = min(current_position + chunk_size, len(text))
        
        if chunk_end < len(text):
            # Find the last sentence break within the chunk
            last_break = max(
                text.rfind('. ', current_position, chunk_end),
                text.rfind('? ', current_position, chunk_end),
                text.rfind('! ', current_position, chunk_end),
                text.rfind('\n', current_position, chunk_end)
            )
            
            if last_break != -1:
                chunk_end = last_break + 1
        
        chunk = text[current_position:chunk_end].strip()
        if chunk:
            chunks.append(chunk)
        
        current_position = chunk_end - chunk_overlap
        if current_position < 0 or current_position >= len(text):
            break
    
    return chunks

def load_and_chunk_documents(directory: str) -> List[Dict[str, Any]]:
    """Load and chunk documents from a directory"""
    chunks = []
    supported_extensions = {
        '.txt': process_text_file,
        '.md': process_text_file,
        '.csv': process_text_file,
        '.html': process_text_file,
        '.pdf': process_pdf_file,
        '.docx': process_docx_file,
        '.doc': process_docx_file,
        '.xlsx': process_excel_file,
        '.xls': process_excel_file,
        '.epub': process_epub_file,
        '.pptx': process_pptx_file,
        '.rtf': process_rtf_file,
        '.json': process_json_file,
        '.xml': process_xml_file
    }
    
    # Convert all extensions to lowercase for case-insensitive comparison
    image_extensions = {ext.lower() for ext in IMAGE_EXTENSIONS}
    supported_extensions = {ext.lower(): func for ext, func in supported_extensions.items()}
    
    for filename in os.listdir(directory):
        filepath = os.path.join(directory, filename)
        if os.path.isfile(filepath):
            file_ext = Path(filepath).suffix.lower()  # Convert to lowercase
            
            # Skip image files with a single message
            if file_ext in image_extensions:
                print(f"Skipping image file: {filename}")
                continue
            
            if file_ext in supported_extensions:
                try:
                    process_func = supported_extensions[file_ext]
                    file_chunks = process_func(filepath)
                    
                    # Add non-empty chunks with source information
                    valid_chunks = []
                    for chunk in file_chunks:
                        if chunk and len(chunk.strip()) > 20:
                            chunk_with_source = f"{chunk.strip()} [Source: {filename}]"
                            valid_chunks.append(chunk_with_source)
                    
                    chunks.extend(valid_chunks)
                    if valid_chunks:
                        print(f"Processed: {filename} - Found {len(valid_chunks)} chunks.")
                    else:
                        print(f"Processed: {filename} - No valid chunks found (file might be empty or contain only small text fragments).")
                except Exception as e:
                    print(f"Error processing {filename}: {e}")
                    if file_ext == '.pptx' and not PPTX_AVAILABLE:
                        print(f"Note: PPTX processing requires python-pptx package. Install with: pip install python-pptx")
            else:
                print(f"Skipping unsupported file type: {filename}")
    
    return chunks

class DocumentHandler(FileSystemEventHandler):
    def __init__(self, documents_dir: str):
        self.documents_dir = documents_dir
        self.processed_files = set()
        self.load_processed_files()

    def load_processed_files(self):
        """Load the list of already processed files"""
        try:
            with open('processed_files.json', 'r') as f:
                self.processed_files = set(json.load(f))
        except FileNotFoundError:
            self.processed_files = set()

    def save_processed_files(self):
        """Save the list of processed files"""
        with open('processed_files.json', 'w') as f:
            json.dump(list(self.processed_files), f)

    def on_created(self, event):
        """Handle new file creation"""
        if not event.is_directory and event.src_path.endswith(('.txt', '.pdf', '.docx', '.xlsx', '.epub', '.pptx', '.rtf', '.json', '.xml')):
            file_path = event.src_path
            if file_path not in self.processed_files:
                print(f"New file detected: {file_path}")
                self.process_file(file_path)
                self.processed_files.add(file_path)
                self.save_processed_files()
                # Trigger generate_embeddings.py
                self.generate_embeddings()

    def process_file(self, file_path: str):
        """Process a single file"""
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            chunks = []

            if file_extension == '.txt':
                chunks = process_text_file(file_path)
            elif file_extension == '.pdf':
                chunks = process_pdf_file(file_path)
            elif file_extension == '.docx':
                chunks = process_docx_file(file_path)
            elif file_extension in ['.xlsx', '.xls']:
                chunks = process_excel_file(file_path)
            elif file_extension == '.epub':
                chunks = process_epub_file(file_path)
            elif file_extension == '.pptx':
                chunks = process_pptx_file(file_path)
            elif file_extension == '.rtf':
                chunks = process_rtf_file(file_path)
            elif file_extension == '.json':
                chunks = process_json_file(file_path)
            elif file_extension == '.xml':
                chunks = process_xml_file(file_path)

            if chunks:
                # Save chunks to temp_chunks.txt
                with open('temp_chunks.txt', 'a', encoding='utf-8') as f:
                    for chunk in chunks:
                        f.write(f"{file_path}|{chunk}\n")
                print(f"Processed {len(chunks)} chunks from {file_path}")

        except Exception as e:
            print(f"Error processing file {file_path}: {e}")

    def generate_embeddings(self):
        """Trigger the generate_embeddings.py script"""
        try:
            subprocess.run([sys.executable, 'scripts/generate_embeddings.py'], check=True)
            print("Embeddings generated successfully")
        except subprocess.CalledProcessError as e:
            print(f"Error generating embeddings: {e}")

def watch_directory(documents_dir: str):
    """Watch the documents directory for new files"""
    event_handler = DocumentHandler(documents_dir)
    observer = Observer()
    observer.schedule(event_handler, documents_dir, recursive=False)
    observer.start()
    
    print(f"Watching directory: {documents_dir}")
    print("Press Ctrl+C to stop watching")
    
    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
        print("\nStopping file watcher...")
    
    observer.join()

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python load_and_chunk.py <documents_directory>")
        sys.exit(1)

    documents_dir = sys.argv[1]
    if not os.path.isdir(documents_dir):
        print(f"Error: {documents_dir} is not a valid directory")
        sys.exit(1)

    # Start watching the directory
    watch_directory(documents_dir) 